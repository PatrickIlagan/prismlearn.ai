"""
Text extraction for the four supported source formats.

Each extractor enforces the PRD hard limits (Doc 1 §1) and returns a normalized
ExtractedSource. Extraction is deliberately dumb — it only pulls text; structuring
into chapters is the LLM's job in [MODE: INGEST].
"""

from __future__ import annotations

import io
import ipaddress
import re
import socket
from urllib.parse import parse_qs, urlparse

import httpx
import pdfplumber
import trafilatura
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.core.config import settings
from app.schemas.ingest import ExtractedSource, SourceType
from app.services import vision


class ExtractionError(Exception):
    """Raised for unprocessable sources (too long, unreadable, no text)."""


# --------------------------------------------------------------------------- #
# PPTX
# --------------------------------------------------------------------------- #
def extract_pptx(data: bytes, filename: str) -> ExtractedSource:
    try:
        prs = Presentation(io.BytesIO(data))
    except Exception as exc:  # noqa: BLE001 - surface a clean 4xx to the client
        raise ExtractionError(f"Could not open the presentation: {exc}") from exc

    slides = prs.slides
    if len(slides) > settings.max_document_pages:
        raise ExtractionError(
            f"Presentation has {len(slides)} slides; the limit is {settings.max_document_pages}."
        )

    # Gather text and embedded pictures per slide first, so captioning (one
    # Gemma call per image, when enabled) happens as a single capped batch
    # across the whole deck rather than per-slide. When vision.captions_enabled()
    # is False (the default), caption_images() is a no-op and this whole path
    # produces byte-identical output to before slide images were considered.
    slide_text: dict[int, list[str]] = {}
    slide_images: list[tuple[int, bytes]] = []
    for idx, slide in enumerate(slides, start=1):
        lines = [
            shape.text.strip()
            for shape in slide.shapes
            if shape.has_text_frame and shape.text.strip()
        ]
        if lines:
            slide_text[idx] = lines
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    slide_images.append((idx, shape.image.blob))
                except Exception:  # noqa: BLE001 - unsupported/corrupt image, skip it
                    continue

    captions = vision.caption_images([blob for _, blob in slide_images])
    captions_by_slide: dict[int, list[str]] = {}
    for (idx, _blob), caption in zip(slide_images, captions):
        if caption:
            captions_by_slide.setdefault(idx, []).append(caption)

    chunks: list[str] = []
    for idx in range(1, len(slides) + 1):
        body = list(slide_text.get(idx, []))
        body.extend(f"[Figure: {c}]" for c in captions_by_slide.get(idx, []))
        if body:
            chunks.append(f"[Slide {idx}]\n" + "\n".join(body))

    text = "\n\n".join(chunks).strip()
    if not text:
        raise ExtractionError("No extractable text found in the presentation.")

    return ExtractedSource(
        source_type=SourceType.pptx,
        title=_clean_title(filename),
        raw_text=text,
        unit_count=len(slides),
    )


# --------------------------------------------------------------------------- #
# PDF
# --------------------------------------------------------------------------- #
def extract_pdf(data: bytes, filename: str) -> ExtractedSource:
    try:
        pdf = pdfplumber.open(io.BytesIO(data))
    except Exception as exc:  # noqa: BLE001
        raise ExtractionError(f"Could not open the PDF: {exc}") from exc

    with pdf:
        if len(pdf.pages) > settings.max_document_pages:
            raise ExtractionError(
                f"PDF has {len(pdf.pages)} pages; the limit is {settings.max_document_pages}."
            )
        chunks: list[str] = []
        for idx, page in enumerate(pdf.pages, start=1):
            page_text = (page.extract_text() or "").strip()
            if page_text:
                chunks.append(f"[Page {idx}]\n{page_text}")
        page_count = len(pdf.pages)

    text = "\n\n".join(chunks).strip()
    if not text:
        raise ExtractionError(
            "No selectable text found in the PDF (it may be a scanned image requiring OCR)."
        )

    return ExtractedSource(
        source_type=SourceType.pdf,
        title=_clean_title(filename),
        raw_text=text,
        unit_count=page_count,
    )


# --------------------------------------------------------------------------- #
# YouTube
# --------------------------------------------------------------------------- #
def extract_youtube(url: str) -> ExtractedSource:
    # youtube-transcript-api v1.x is instance-based (.fetch), not the old
    # classmethod .get_transcript. Imported lazily to keep import cost off startup.
    from youtube_transcript_api import (
        YouTubeTranscriptApi,
        NoTranscriptFound,
        TranscriptsDisabled,
        VideoUnavailable,
    )

    video_id = _parse_youtube_id(url)
    if not video_id:
        raise ExtractionError("Could not parse a YouTube video ID from that URL.")

    try:
        fetched = YouTubeTranscriptApi().fetch(video_id)
        segments = fetched.to_raw_data()  # -> [{"text", "start", "duration"}, ...]
    except (TranscriptsDisabled, NoTranscriptFound):
        raise ExtractionError("This video has no available transcript/captions.")
    except VideoUnavailable:
        raise ExtractionError("This video is unavailable.")
    except Exception as exc:  # noqa: BLE001
        raise ExtractionError(f"Failed to fetch transcript: {exc}") from exc

    if not segments:
        raise ExtractionError("The transcript for this video is empty.")

    # Duration = start of last segment + its duration.
    last = segments[-1]
    total_seconds = float(last.get("start", 0)) + float(last.get("duration", 0))
    minutes = total_seconds / 60.0
    if minutes > settings.max_youtube_minutes:
        raise ExtractionError(
            f"Video is ~{minutes:.0f} min long; the limit is {settings.max_youtube_minutes} min."
        )

    text = " ".join(seg["text"].strip() for seg in segments if seg.get("text", "").strip())
    return ExtractedSource(
        source_type=SourceType.youtube,
        title=f"YouTube video {video_id}",
        raw_text=text,
        unit_count=max(1, round(minutes)),
    )


# --------------------------------------------------------------------------- #
# Website
# --------------------------------------------------------------------------- #
_WEBSITE_MAX_BYTES = 3 * 1024 * 1024  # 3 MB of raw HTML is plenty for an article
_WEBSITE_UA = {"User-Agent": "Mozilla/5.0 (compatible; PrismLearnBot/1.0)"}


def is_youtube_url(url: str) -> bool:
    """Public check the /ingest/url router uses to pick an extractor."""
    return _parse_youtube_id(url) is not None


def _is_safe_host(hostname: str) -> bool:
    """Basic SSRF guard: resolve the hostname and reject private/loopback/
    link-local/reserved addresses (this also catches cloud metadata endpoints
    like 169.254.169.254, which is link-local) before fetching a
    user-supplied URL server-side. Best-effort, not exhaustive DNS-rebinding
    protection — good enough for a hackathon-scale deployment.
    """
    if not hostname:
        return False
    try:
        infos = socket.getaddrinfo(hostname, None)
    except socket.gaierror:
        return False
    for info in infos:
        ip = ipaddress.ip_address(info[4][0])
        if ip.is_private or ip.is_loopback or ip.is_link_local or ip.is_reserved or ip.is_multicast:
            return False
    return True


def _fetch_html(url: str) -> bytes:
    """Bounded fetch with a manual, re-validated redirect chain — each hop's
    host is checked (not just the original URL), so a public URL that
    redirects to an internal address is still caught. Returns raw bytes
    (not decoded) — trafilatura does its own encoding detection from bytes,
    and does it more reliably than a naive UTF-8 decode here would."""
    current = url.strip()
    for _ in range(5):
        parsed = urlparse(current)
        if parsed.scheme not in ("http", "https") or not parsed.hostname:
            raise ExtractionError("That doesn't look like a valid website URL.")
        if not _is_safe_host(parsed.hostname):
            raise ExtractionError("That URL points at a private/internal address, which isn't allowed.")

        try:
            with httpx.stream(
                "GET", current, follow_redirects=False, timeout=8.0, headers=_WEBSITE_UA
            ) as resp:
                if resp.is_redirect:
                    location = resp.headers.get("location")
                    if not location:
                        raise ExtractionError("That URL redirected without a destination.")
                    current = str(httpx.URL(current).join(location))
                    continue
                resp.raise_for_status()
                content_type = resp.headers.get("content-type", "")
                if content_type and "html" not in content_type:
                    raise ExtractionError(
                        f"That URL doesn't look like a webpage (content-type: {content_type})."
                    )
                buf = bytearray()
                for chunk in resp.iter_bytes():
                    buf.extend(chunk)
                    if len(buf) > _WEBSITE_MAX_BYTES:
                        raise ExtractionError("That page is too large to ingest.")
                return bytes(buf)
        except httpx.HTTPStatusError as exc:
            raise ExtractionError(f"That URL returned an error ({exc.response.status_code}).") from exc
        except httpx.TimeoutException as exc:
            # Corporate/marketing sites are frequently behind bot-detection (Akamai,
            # Cloudflare, etc.) that silently drops or stalls automated requests
            # rather than returning a clean block page — from the caller's side
            # this is indistinguishable from the site just being slow. Confirmed
            # against amd.com/en/blogs/... via two independent network paths (this
            # server, and a completely separate fetch service) — both failed the
            # same way, so it's the target site, not this code.
            raise ExtractionError(
                "That page didn't respond in time. Some sites block automated "
                "requests like this one — if it's a well-known page, try pasting "
                "the article text directly, or a different source."
            ) from exc
        except httpx.HTTPError as exc:
            raise ExtractionError(f"Could not reach that URL: {exc}") from exc
    raise ExtractionError("Too many redirects.")


def extract_website(url: str) -> ExtractedSource:
    html = _fetch_html(url)

    text = trafilatura.extract(html, favor_recall=True)
    if not text or not text.strip():
        raise ExtractionError("Could not find readable article text on that page.")

    if len(text) > settings.max_website_chars:
        raise ExtractionError(
            f"That page has ~{len(text)} characters of text; the limit is {settings.max_website_chars}."
        )

    metadata = trafilatura.extract_metadata(html)
    title = ((metadata.title if metadata else None) or urlparse(url).hostname or "Untitled").strip()

    return ExtractedSource(
        source_type=SourceType.website,
        title=title,
        raw_text=text,
        unit_count=max(1, round(len(text) / 1000)),
    )


# --------------------------------------------------------------------------- #
# Helpers
# --------------------------------------------------------------------------- #
def _parse_youtube_id(url: str) -> str | None:
    parsed = urlparse(url.strip())
    host = (parsed.hostname or "").lower().removeprefix("www.")
    if host == "youtu.be":
        return parsed.path.lstrip("/") or None
    if host in {"youtube.com", "m.youtube.com", "music.youtube.com"}:
        if parsed.path == "/watch":
            return parse_qs(parsed.query).get("v", [None])[0]
        for prefix in ("/embed/", "/shorts/", "/v/"):
            if parsed.path.startswith(prefix):
                return parsed.path[len(prefix):].split("/")[0] or None
    # Bare 11-char id fallback.
    if re.fullmatch(r"[A-Za-z0-9_-]{11}", url.strip()):
        return url.strip()
    return None


def _clean_title(filename: str) -> str:
    stem = re.sub(r"\.(pdf|pptx)$", "", filename, flags=re.IGNORECASE)
    stem = re.sub(r"[_-]+", " ", stem).strip()
    return stem or "Untitled"
